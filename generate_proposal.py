import os
from dotenv import load_dotenv
from langchain_openai import ChatOpenAI
from langchain.prompts import ChatPromptTemplate
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime

# Load API key
load_dotenv()

def read_input():
    """Read requirements from input.txt file"""
    try:
        with open("input.txt", "r", encoding="utf-8") as f:
            return f.read().strip()
    except FileNotFoundError:
        print("❌ File input.txt not found")
        print("💡 Please create input.txt and paste client requirements there")
        return None

def generate_content_with_ai(requirement_text):
    """Send prompt to GPT-4 to generate content for 5 slides"""
    
    llm = ChatOpenAI(
        model="gpt-4o-mini",  # Using mini for cost efficiency, can switch to "gpt-4o" if needed
        temperature=0.7
    )
    
    prompt = ChatPromptTemplate.from_messages([
        ("system", """You are a professional Solution Architect at a technology company.
Writing style: professional, concise, technical-focused, avoid clichés.
Please write content for 5 proposal slides based on the project requirements.

Output format: Each slide separated by "---SLIDE---"
Slide 1: INTRODUCTION - Company introduction and capabilities
Slide 2: PROBLEM STATEMENT - Analysis of client's current challenges
Slide 3: SOLUTION - Specific proposed solution
Slide 4: TECHNOLOGY STACK - Technologies to be used
Slide 5: TIMELINE - Project implementation roadmap

Each slide should have 3-5 bullet points, concise and easy to understand."""),
        ("user", "Project requirements:\n\n{requirement}")
    ])
    
    print("🤖 Sending request to AI...")
    chain = prompt | llm
    response = chain.invoke({"requirement": requirement_text})
    
    return response.content

def parse_slides(ai_content):
    """Parse AI response content into 5 separate slides"""
    slides = ai_content.split("---SLIDE---")
    
    # Clean and format
    parsed = []
    for i, slide in enumerate(slides, 1):
        slide = slide.strip()
        # Remove header if exists (e.g., "Slide 1: INTRODUCTION")
        lines = slide.split('\n')
        content_lines = []
        for line in lines:
            line = line.strip()
            if line and not line.startswith("Slide"):
                content_lines.append(line)
        
        parsed.append('\n'.join(content_lines))
    
    # Ensure we have 5 slides
    while len(parsed) < 5:
        parsed.append("Content is being updated...")
    
    return parsed[:5]

def create_powerpoint(slides_content):
    """Create PowerPoint file from 5 slides content"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide titles
    slide_titles = [
        "INTRODUCTION",
        "PROBLEM STATEMENT",
        "SOLUTION",
        "TECHNOLOGY STACK",
        "TIMELINE"
    ]
    
    for i, (title, content) in enumerate(zip(slide_titles, slides_content)):
        # Use Title and Content layout
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Content
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.word_wrap = True
        
        # Parse content into bullet points
        lines = content.split('\n')
        for idx, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
            
            # Remove bullet characters if any
            line = line.lstrip('•-*').strip()
            
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = line
            p.font.size = Pt(18)
            p.level = 0
            p.space_after = Pt(12)
    
    # Create filename with timestamp
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_filename = f"proposal_{timestamp}.pptx"
    
    prs.save(output_filename)
    return output_filename

def main():
    """Main function - CLI entry point"""
    print("=" * 50)
    print("🚀 BIDCRAFT - Proposal Generator")
    print("=" * 50)
    print()
    
    # Step 1: Read input
    requirement = read_input()
    if not requirement:
        return
    
    print(f"📄 Read requirements ({len(requirement)} characters)")
    print()
    
    # Step 2: Call AI
    try:
        ai_content = generate_content_with_ai(requirement)
        print("✅ Received response from AI")
        print()
    except Exception as e:
        print(f"❌ Error calling AI: {e}")
        print("💡 Please check OPENAI_API_KEY in .env file")
        return
    
    # Step 3: Parse slides
    slides = parse_slides(ai_content)
    print(f"📊 Parsed into {len(slides)} slides")
    print()
    
    # Step 4: Create PowerPoint
    try:
        output_file = create_powerpoint(slides)
        print("=" * 50)
        print(f"✅ COMPLETED!")
        print(f"📁 Output file: {output_file}")
        print("=" * 50)
    except Exception as e:
        print(f"❌ Error creating PowerPoint: {e}")
        return

if __name__ == "__main__":
    main()


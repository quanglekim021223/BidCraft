"""PowerPoint generation service"""

from datetime import datetime
from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt

from app.config.settings import Settings
from app.utils.file_utils import ensure_output_dir


class PPTXService:
    """Service for creating PowerPoint presentations"""
    
    def __init__(self):
        """Initialize PowerPoint service"""
        self.slide_titles = Settings.SLIDE_TITLES
        self.slide_width = Inches(Settings.SLIDE_WIDTH)
        self.slide_height = Inches(Settings.SLIDE_HEIGHT)
        self.title_font_size = Pt(Settings.TITLE_FONT_SIZE)
        self.content_font_size = Pt(Settings.CONTENT_FONT_SIZE)
    
    def create_presentation(
        self, 
        slides_content: List[str], 
        output_dir: str = None
    ) -> str:
        """
        Create PowerPoint file from slides content
        
        Args:
            slides_content: List of content for each slide
            output_dir: Output directory (default: from settings)
            
        Returns:
            Path to created PowerPoint file
        """
        output_dir = output_dir or Settings.OUTPUT_DIR
        ensure_output_dir(output_dir)
        
        prs = Presentation()
        prs.slide_width = self.slide_width
        prs.slide_height = self.slide_height
        
        # Create slides
        for title, content in zip(self.slide_titles, slides_content):
            self._add_slide(prs, title, content)
        
        # Generate filename with timestamp
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_filename = f"proposal_{timestamp}.pptx"
        output_path = Path(output_dir) / output_filename
        
        prs.save(str(output_path))
        return str(output_path)
    
    def _add_slide(self, presentation: Presentation, title: str, content: str):
        """
        Add a slide to the presentation
        
        Args:
            presentation: PowerPoint presentation object
            title: Slide title
            content: Slide content
        """
        # Use Title and Content layout
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = self.title_font_size
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Set content
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.word_wrap = True
        
        # Parse content into bullet points
        lines = content.split('\n')
        for idx, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
            
            # Remove bullet characters if any
            line = line.lstrip('•-*').strip()
            
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = line
            p.font.size = self.content_font_size
            p.level = 0
            p.space_after = Pt(12)

